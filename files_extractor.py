from pptx import Presentation
import os
import GPT
from PyPDF2 import PdfReader


def get_ppt_text(pptx_filename):
    text=""

    with open(f"{pptx_filename}","rb") as f:
        prs = Presentation(f)
        for slide in prs.slides:
            for shape in slide.shapes:
                if  shape.has_text_frame:
                    text+=str(shape.text)
                    print(text)
    return text

def get_pdf_text(pdf_filename):
    text=""
    reader = PdfReader(pdf_filename)
    for page in reader.pages:
        text+=str(page.extract_text())
        print(text)
    return text

#This function will filter the text and according to the given keywords
def find_content(keywords,text):
    content = []
    with open("content.txt",'w',encoding="utf-8") as f:
        f.writelines(text)
    lines=[]
    with open("content.txt",'r',encoding="utf-8") as file:
        lines=file.readlines()
        # print(lines)
    for line in lines:
        for keyword in keywords:
            if keyword.lower() in list(line.split(" ")):
                if line not in content:
                    # print(f"upadated in text-----------------{line}")
                    content.append(line)
            elif keyword.title() in list(line.split(" ")):
                if line not in content:
                    # print(f"upadated in text-----------------{line}")
                    content.append(line)
            elif keyword.upper() in list(line.split(" ")):
                if line not in content:
                    # print(f"upadated in text-----------------{line}")
                    content.append(line)

    return "".join(content)
